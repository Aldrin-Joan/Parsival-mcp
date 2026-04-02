import pytest

try:
    import pptx
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    pptx = None

from src.models.enums import FileFormat, ParseStatus, SectionType
from src.parsers.pptx_parser import PptxParser


@pytest.mark.skipif(pptx is None, reason="python-pptx required")
@pytest.mark.asyncio
async def test_pptx_parser_mixed_content(tmp_path):
    file_path = tmp_path / "mixed.pptx"
    prs = Presentation()

    # Slide 1: title + paragraph + table + image + notes
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    if title is None:
        title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title.text = "Sample Slide Title"

    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.5))
    textbox.text = "This is sample text in slide 1."

    table_shape = slide.shapes.add_table(3, 2, Inches(1), Inches(3), Inches(8), Inches(2)).table
    table_shape.cell(0, 0).text = "H1"
    table_shape.cell(0, 1).text = "H2"
    table_shape.cell(1, 0).text = "A"
    table_shape.cell(1, 1).text = "B"
    table_shape.cell(2, 0).text = "C"
    table_shape.cell(2, 1).text = "D"

    # Add image
    img_path = tmp_path / "img.png"
    from PIL import Image

    img = Image.new("RGB", (10, 10), color="blue")
    img.save(img_path)
    slide.shapes.add_picture(str(img_path), Inches(1), Inches(5), width=Inches(1), height=Inches(1))

    # Speaker notes
    notes = slide.notes_slide
    notes.notes_text_frame.text = "Speaker notes here."

    # Slide 2: empty slide
    prs.slides.add_slide(prs.slide_layouts[6])

    # Slide 3: nested group shape
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    group = slide3.shapes.add_group_shape()
    inner_txt = group.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    inner_txt.text = "Nested text shape"

    prs.save(str(file_path))

    parser = PptxParser()
    result = await parser.parse(file_path)

    assert result.status == ParseStatus.OK
    assert result.metadata.file_format == FileFormat.PPTX
    assert result.metadata.page_count == 3
    assert any(s.type == SectionType.HEADING and "Sample Slide Title" in s.content for s in result.sections)
    assert any(s.type == SectionType.PARAGRAPH and "This is sample text" in s.content for s in result.sections)
    assert result.metadata.table_count == 1
    assert result.metadata.image_count == 1
    assert any(s.type == SectionType.METADATA and "Speaker notes here." in s.content for s in result.sections)


@pytest.mark.skipif(pptx is None, reason="python-pptx required")
@pytest.mark.asyncio
async def test_pptx_parser_empty_slides(tmp_path):
    file_path = tmp_path / "empty.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(str(file_path))

    parser = PptxParser()
    result = await parser.parse(file_path)

    assert result.status == ParseStatus.OK
    assert result.metadata.page_count == 2
    assert len(result.sections) == 0
    assert result.metadata.table_count == 0
    assert result.metadata.image_count == 0


@pytest.mark.skipif(pptx is None, reason="python-pptx required")
@pytest.mark.asyncio
async def test_pptx_parser_nested_shapes(tmp_path):
    file_path = tmp_path / "nested.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    group = slide.shapes.add_group_shape()
    inner = group.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    inner.text = "Nested paragraph"
    prs.save(str(file_path))

    parser = PptxParser()
    result = await parser.parse(file_path)

    assert result.status == ParseStatus.OK
    assert any("Nested paragraph" in s.content for s in result.sections)
