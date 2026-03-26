from __future__ import annotations
import time
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from src.models.enums import FileFormat, ParseStatus, SectionType
from src.models.metadata import DocumentMetadata
from src.models.parse_result import ParseResult, ParseError, Section
from src.models.table import TableResult, TableCell
from src.models.image import ImageRef
from src.parsers.base import BaseParser
from src.parsers.registry import register


def _shape_text(shape) -> str:
    if not shape.has_text_frame:
        return ''
    texts = []
    for paragraph in shape.text_frame.paragraphs:
        text = paragraph.text.strip()
        if text:
            texts.append(text)
    return '\n'.join(texts)


def _extract_table(shape, slide_number: int, index: int) -> TableResult:
    table = shape.table
    rows_data = []
    cells = []
    for r_idx, row in enumerate(table.rows):
        row_values = []
        for c_idx, cell in enumerate(row.cells):
            value = (cell.text or '').strip()
            row_values.append(value)
            cells.append(
                TableCell(
                    row=r_idx,
                    col=c_idx,
                    value=value,
                    raw_value=value,
                    colspan=1,
                    rowspan=1,
                    is_header=(r_idx == 0),
                )
            )
        rows_data.append(row_values)

    headers = rows_data[0] if rows_data else []
    body_rows = rows_data[1:] if len(rows_data) > 1 else []

    return TableResult(
        index=index,
        page=slide_number,
        caption=None,
        headers=headers,
        rows=body_rows,
        cells=cells,
        row_count=max(0, len(body_rows)),
        col_count=len(headers),
        has_merged_cells=False,
        confidence=0.9,
        confidence_reason='pptx parser',
        markdown='',
        errors=[],
    )


def _extract_image(shape, slide_number: int, index: int) -> ImageRef:
    image = shape.image
    blob = image.blob
    fmt = image.content_type.split('/')[-1] if image.content_type else 'png'
    b64 = __import__('base64').b64encode(blob).decode('ascii')
    return ImageRef(
        index=index,
        page=slide_number,
        width_px=image.size[0],
        height_px=image.size[1],
        format=fmt,
        size_bytes=len(blob),
        base64_data=b64,
        description_hint=f'PPTX slide {slide_number} image {index}',
        confidence=1.0,
        alt_text=getattr(shape, 'name', None),
    )


def _walk_shapes(shapes, slide_number: int, table_start: int, image_start: int, section_start: int):
    sections: list[Section] = []
    tables: list[TableResult] = []
    images: list[ImageRef] = []
    section_idx = section_start
    table_idx = table_start
    image_idx = image_start

    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            s, t, i, section_idx, table_idx, image_idx = _walk_shapes(
                shape.shapes, slide_number, table_idx, image_idx, section_idx
            )
            sections.extend(s)
            tables.extend(t)
            images.extend(i)
            continue

        if shape.has_table:
            table = _extract_table(shape, slide_number, table_idx)
            tables.append(table)
            sections.append(
                Section(
                    index=section_idx,
                    type=SectionType.TABLE,
                    content='',
                    page=slide_number,
                    level=None,
                    table=table,
                    confidence=table.confidence,
                )
            )
            table_idx += 1
            section_idx += 1
            continue

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img = _extract_image(shape, slide_number, image_idx)
            images.append(img)
            sections.append(
                Section(
                    index=section_idx,
                    type=SectionType.IMAGE,
                    content='',
                    page=slide_number,
                    level=None,
                    images=[img],
                    confidence=img.confidence,
                )
            )
            image_idx += 1
            section_idx += 1
            continue

        text = _shape_text(shape)
        if text:
            stype = SectionType.PARAGRAPH
            if getattr(shape, 'is_placeholder', False) and 'title' in (shape.name or '').lower():
                stype = SectionType.HEADING
            sections.append(
                Section(
                    index=section_idx,
                    type=stype,
                    content=text,
                    page=slide_number,
                    level=1 if stype == SectionType.HEADING else None,
                    confidence=0.95 if stype == SectionType.HEADING else 0.85,
                )
            )
            section_idx += 1

    return sections, tables, images, section_idx, table_idx, image_idx


@register(FileFormat.PPTX)
class PptxParser(BaseParser):

    async def parse(self, path: Path, options: dict | None = None) -> ParseResult:
        src = Path(path)
        if not src.exists():
            return ParseResult(
                status=ParseStatus.FAILED,
                metadata=DocumentMetadata(source_path=str(src), file_format=FileFormat.PPTX, file_size_bytes=0),
                sections=[],
                images=[],
                tables=[],
                errors=[ParseError(code='not_found', message='PPTX file not found', recoverable=False)],
                raw_text='',
                cache_hit=False,
                request_id='',
            )

        start = time.time()
        prs = Presentation(str(src))

        sections: list[Section] = []
        tables: list[TableResult] = []
        images: list[ImageRef] = []
        errors: list[ParseError] = []

        section_idx = 0
        table_idx = 0
        image_idx = 0

        for slide_number, slide in enumerate(prs.slides, start=1):
            title_text = None
            if slide.shapes.title is not None and slide.shapes.title.text.strip():
                title_text = slide.shapes.title.text.strip()

            if title_text:
                sections.append(
                    Section(
                        index=section_idx,
                        type=SectionType.HEADING,
                        content=title_text,
                        page=slide_number,
                        level=1,
                        confidence=1.0,
                    )
                )
                section_idx += 1

            s, t, i, section_idx, table_idx, image_idx = _walk_shapes(
                slide.shapes, slide_number, table_idx, image_idx, section_idx
            )
            sections.extend(s)
            tables.extend(t)
            images.extend(i)

            if hasattr(slide, 'notes_slide') and slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    sections.append(
                        Section(
                            index=section_idx,
                            type=SectionType.METADATA,
                            content=notes,
                            page=slide_number,
                            level=None,
                            confidence=0.7,
                        )
                    )
                    section_idx += 1

        metadata = DocumentMetadata(
            source_path=str(src),
            file_format=FileFormat.PPTX,
            file_size_bytes=src.stat().st_size,
            title=prs.core_properties.title,
            author=prs.core_properties.author,
            subject=prs.core_properties.subject,
            keywords=(prs.core_properties.keywords.split(',') if prs.core_properties.keywords else []),
            page_count=len(prs.slides),
            section_count=len(sections),
            table_count=len(tables),
            image_count=len(images),
            has_toc=any(s.type == SectionType.HEADING for s in sections),
            toc=[],
            word_count=sum(len(s.content.split()) for s in sections),
            char_count=sum(len(s.content) for s in sections),
            reading_time_minutes=None,
            parse_duration_ms=(time.time() - start) * 1000,
            parser_version='pptx',
        )

        return ParseResult(
            status=ParseStatus.OK,
            metadata=metadata,
            sections=sections,
            images=images,
            tables=tables,
            errors=errors,
            raw_text='\n'.join(s.content for s in sections if s.content),
            cache_hit=False,
            request_id='',
        )

    async def parse_metadata(self, path: Path) -> DocumentMetadata:
        src = Path(path)
        if not src.exists():
            raise FileNotFoundError('PPTX file not found')

        prs = Presentation(str(src))
        return DocumentMetadata(
            source_path=str(src),
            file_format=FileFormat.PPTX,
            file_size_bytes=src.stat().st_size,
            title=prs.core_properties.title,
            author=prs.core_properties.author,
            subject=prs.core_properties.subject,
            keywords=(prs.core_properties.keywords.split(',') if prs.core_properties.keywords else []),
            page_count=len(prs.slides),
            section_count=0,
            table_count=0,
            image_count=0,
            has_toc=False,
            toc=[],
            parse_duration_ms=0.0,
            parser_version='pptx',
        )
